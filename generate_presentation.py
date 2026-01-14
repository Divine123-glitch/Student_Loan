"""
Generate NELFUND Navigator PowerPoint Presentation
Run: python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
PURPLE_DARK = RGBColor(102, 126, 234)  # #667eea
PURPLE_LIGHT = RGBColor(118, 75, 162)  # #764ba2
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(85, 85, 85)
GREEN = RGBColor(16, 185, 129)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(220, 220, 220)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = PURPLE_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)
    
    return slide

# Slide 1: Title
add_title_slide(prs, "🎓 NELFUND Navigator", 
                "AI-Powered Student Loan Assistant for Nigerian Students")

# Slide 2: The Problem
add_content_slide(prs, "🎯 The Problem", [
    "❌ Information Overload - 500+ pages of policy documents scattered across PDFs",
    "❌ Confusion & Uncertainty - 'Am I eligible?' 'What documents do I need?'",
    "❌ Misinformation - Social media rumors leading to wrong decisions",
    "❌ Complex Language - Legal terminology that confuses students",
    "💡 Students need simple, accurate answers - not 500-page PDFs"
])

# Slide 3: Our Solution
add_content_slide(prs, "💡 Our Solution", [
    "🤖 Agentic RAG System - Smart document retrieval with conditional logic",
    "💬 Natural Conversations - Ask in plain English, get document-backed answers",
    "📚 Source Citations - Every answer includes references to official documents",
    "🧠 Conversation Memory - Remembers context for intelligent follow-ups",
    "🎯 Your Path to Higher Education Starts Here"
])

# Slide 4: Tech Stack
add_two_column_slide(prs, "🛠️ Technology Stack",
    ["Backend:", "• FastAPI", "• LangChain", "• LangGraph", "• OpenAI GPT-4", "• ChromaDB", "• JWT Auth"],
    ["Frontend:", "• React 18", "• Vite", "• Tailwind CSS", "• React Router", "• Axios"]
)

# Slide 5: Architecture
add_content_slide(prs, "🏗️ System Architecture", [
    "User Question → Frontend (React) → Backend API (FastAPI)",
    "→ Agent Classification (LangGraph) → Document Retrieval (Conditional)",
    "→ LLM Generation (GPT-4) → Response + Sources",
    "🗄️ Dual Database: chroma_users/ (user data) & chroma_db/ (documents)",
    "🔒 Security: JWT tokens, bcrypt passwords, CORS protection"
])

# Slide 6: Features Overview
add_content_slide(prs, "✨ Key Features", [
    "🏠 Interactive Homepage - Modern design with feature showcase",
    "🔐 Secure Authentication - JWT tokens + bcrypt password hashing",
    "💬 Chat Interface - Claude AI-inspired design",
    "🌙 Dark/Light Mode - Theme toggle for comfortable use",
    "📱 Mobile Responsive - Works perfectly on all devices"
])

# Slide 7: Chat Features
add_two_column_slide(prs, "💬 Chat Interface Features",
    ["User Features:", "• Collapsible sidebar", "• Chat history", "• Suggested prompts", "• Source citations", "• Session management"],
    ["UX Elements:", "• Typing indicators", "• Message bubbles", "• Auto-scroll", "• Loading states", "• Error handling"]
)

# Slide 8: Agentic RAG Magic
add_content_slide(prs, "🤖 The Magic: Agentic Behavior", [
    "✓ Example 1: 'Hello' → No retrieval needed → Quick response",
    "✓ Example 2: 'Am I eligible?' → Retrieve docs → Detailed answer with sources",
    "✓ Example 3: 'What documents?' → Retrieve + Use context → Contextual response",
    "🎯 This saves API costs and provides faster responses!",
    "💡 The system THINKS before acting, not just blindly retrieving"
])

# Slide 9: API Endpoints
add_two_column_slide(prs, "📡 RESTful API",
    ["Authentication:", "• POST /api/auth/register", "• POST /api/auth/login", "• GET /api/auth/me"],
    ["Chat Operations:", "• POST /api/chat", "• GET /api/chat/history", "• GET /api/chat/sessions", "• DELETE /api/chat/session/{id}"]
)

# Slide 10: Data Processing
add_content_slide(prs, "📊 Data & Processing", [
    "✓ 9 NELFUND PDF documents processed",
    "✓ 45 total document pages",
    "✓ 68 optimized chunks for retrieval",
    "✓ 44,645 total characters processed",
    "✓ OpenAI text-embedding-3-small for vector embeddings"
])

# Slide 11: User Experience
add_content_slide(prs, "🎨 User Experience Design", [
    "📱 Mobile-First Approach - All features work on mobile",
    "🌙 Dark Mode - Reduced eye strain for late-night studying",
    "💡 Smart Suggestions - Prompts to guide users",
    "🔐 Per-User Data - Each student's chats are completely private",
    "⚡ Fast Responses - Optimized queries and caching"
])

# Slide 12: Authentication System
add_two_column_slide(prs, "🔐 Secure Authentication",
    ["Registration:", "• Email validation", "• Password hashing (bcrypt)", "• User data storage", "• Account creation"],
    ["Login & Sessions:", "• Email/password auth", "• JWT token generation", "• Auto-redirect to chat", "• Session persistence"]
)

# Slide 13: Key Achievements
add_content_slide(prs, "🏆 Key Achievements", [
    "✅ Agentic RAG system fully functional with conditional logic",
    "✅ Full-stack implementation (Frontend + Backend + Database)",
    "✅ 8 RESTful API endpoints with proper authentication",
    "✅ Per-user chat storage and history retrieval",
    "✅ Production-ready code with error handling"
])

# Slide 14: Real Impact
add_two_column_slide(prs, "🌍 Real-World Impact",
    ["By The Numbers:", "• 45 PDF pages processed", "• 68 document chunks", "• 9 NELFUND FAQs covered", "• 24/7 availability"],
    ["Student Benefits:", "• Quick, accurate answers", "• Reduced confusion", "• Better access to info", "• Higher success rate"]
)

# Slide 15: Thank You
add_title_slide(prs, "Thank You! 🎓", 
                "NELFUND Navigator - Empowering Nigerian Students Through AI")

# Save presentation
output_file = "NELFUND_Navigator_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully!")
print(f"📄 File: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")